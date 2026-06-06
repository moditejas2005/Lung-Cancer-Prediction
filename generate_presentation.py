import pptx
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def apply_background(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title(slide, text, subtitle_text=None, top_inch=0.5, height_inch=1.0):
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(top_inch), Inches(11.83), Inches(height_inch))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = 'Calibri'
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(31, 78, 121) # Deep Navy Blue
    
    if subtitle_text:
        p2 = tf.add_paragraph()
        p2.text = subtitle_text
        p2.font.name = 'Calibri'
        p2.font.size = Pt(18)
        p2.font.color.rgb = RGBColor(127, 127, 127) # Muted Gray

def add_bullet_slide(prs, title, bullets):
    slide_layout = prs.slide_layouts[6] # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    apply_background(slide, RGBColor(250, 250, 250))
    
    add_title(slide, title)
    
    # Text Box for Bullets
    txBox = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.83), Inches(5.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    
    for idx, bullet in enumerate(bullets):
        if idx == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
            
        p.space_after = Pt(14)
        
        # Check indentation level (determined by leading tab/spaces or sub-bullets)
        if bullet.startswith("  - "):
            p.level = 1
            p.text = bullet[4:]
            p.font.size = Pt(20)
            p.font.color.rgb = RGBColor(80, 80, 80)
        else:
            p.level = 0
            p.text = bullet.replace("- ", "", 1) if bullet.startswith("- ") else bullet
            p.font.size = Pt(22)
            p.font.bold = True if ":" in p.text.split(" ")[0] else False
            p.font.color.rgb = RGBColor(40, 40, 40)
            
        p.font.name = 'Calibri'
    return slide

def create_presentation():
    prs = pptx.Presentation()
    # 16:9 widescreen layout
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # ─────────────────────────────────────────────────────────────────────────
    # SLIDE 1: Title Slide
    # ─────────────────────────────────────────────────────────────────────────
    slide_layout = prs.slide_layouts[6] # Blank
    slide1 = prs.slides.add_slide(slide_layout)
    apply_background(slide1, RGBColor(245, 248, 250))
    
    # Title Box
    title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.33), Inches(3.0))
    tf = title_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Lung Cancer Prediction Platform"
    p.font.name = 'Calibri'
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(31, 78, 121)
    p.space_after = Pt(8)
    
    p2 = tf.add_paragraph()
    p2.text = "An Explainable and Clinically Validated Calibrated Ensemble Approach"
    p2.font.name = 'Calibri'
    p2.font.size = Pt(22)
    p2.font.color.rgb = RGBColor(46, 117, 182)
    p2.space_after = Pt(24)
    
    p3 = tf.add_paragraph()
    p3.text = "College Demonstration Presentation"
    p3.font.name = 'Calibri'
    p3.font.size = Pt(14)
    p3.font.color.rgb = RGBColor(127, 127, 127)
    
    # ─────────────────────────────────────────────────────────────────────────
    # SLIDE 2: Core Medical Problem & Challenges
    # ─────────────────────────────────────────────────────────────────────────
    add_bullet_slide(prs, "The Core Clinical Problem", [
        "1. Medical Data Scarcity",
        "  - Gathering real-world patient lung cancer records is difficult due to strict privacy rules (HIPAA/GDPR).",
        "2. High Class Imbalance",
        "  - Healthy patient records outnumber positive lung cancer cases, biasing standard model algorithms.",
        "3. The 'Black Box' Barrier",
        "  - Doctors reject standard AI predictions because they cannot explain *why* a patient is labeled high-risk.",
        "4. Overconfident and Uncalibrated Risk Outputs",
        "  - Traditional models give arbitrary probability outputs that do not represent true empirical clinical risk."
    ])
    
    # ─────────────────────────────────────────────────────────────────────────
    # SLIDE 3: Three Pillars of the Platform
    # ─────────────────────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(slide_layout)
    apply_background(slide, RGBColor(250, 250, 250))
    add_title(slide, "System Architecture: Three Pillars")
    
    # 3 Column layout
    col_width = Inches(3.6)
    col_gap = Inches(0.4)
    left_margin = Inches(0.75)
    
    pillars = [
        ("Pillar 1: Data Synthesis", "Conditional Tabular GAN", [
            "Generates highly realistic synthetic patient records.",
            "Learns complex joint probability distributions.",
            "Solves scarcity & maintains privacy."
        ], "D9E1F2"),
        ("Pillar 2: Validation Gates", "Medical & Statistical QA", [
            "Medical rule engine enforces physiological laws.",
            "KS-Test & Chi-Square check data quality.",
            "PSI checks for data drift."
        ], "E2EFDA"),
        ("Pillar 3: Predict & Explain", "Ensemble & SHAP (XAI)", [
            "7-model stacking ensemble for robust risk prediction.",
            "Probability calibration mapping raw scores to absolute risk.",
            "SHAP explains specific risk factors."
        ], "FFF2CC")
    ]
    
    for idx, (title, sub, bullets, bg_color) in enumerate(pillars):
        left = left_margin + idx * (col_width + col_gap)
        
        # Draw background shape
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.8), col_width, Inches(4.8))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor.from_string(bg_color)
        shape.line.color.rgb = RGBColor(200, 200, 200)
        
        # Add Text
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0.2)
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = 'Calibri'
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(31, 78, 121)
        p.space_after = Pt(4)
        
        p2 = tf.add_paragraph()
        p2.text = sub
        p2.font.name = 'Calibri'
        p2.font.size = Pt(14)
        p2.font.italic = True
        p2.font.color.rgb = RGBColor(80, 80, 80)
        p2.space_after = Pt(18)
        
        for b in bullets:
            pb = tf.add_paragraph()
            pb.text = "• " + b
            pb.font.name = 'Calibri'
            pb.font.size = Pt(13)
            pb.font.color.rgb = RGBColor(40, 40, 40)
            pb.space_after = Pt(8)

    # ─────────────────────────────────────────────────────────────────────────
    # SLIDE 4: Pillar 1 - CTGAN Data Synthesis
    # ─────────────────────────────────────────────────────────────────────────
    add_bullet_slide(prs, "Pillar 1: Deep Generative Data Synthesis", [
        "What is CTGAN?",
        "  - Conditional Tabular GAN (Generative Adversarial Network) designed for tabular dataset structures.",
        "Generator vs. Discriminator Competition:",
        "  - The Generator creates fake patient records, trying to fool the Discriminator.",
        "  - The Discriminator learns to distinguish real patient data from synthetic data.",
        "Conditional Masking & Training:",
        "  - CTGAN uses a conditional vector to force the generator to sample imbalanced features (like cancer positive status) evenly.",
        "Rejection Sampling:",
        "  - Generates sub-populations with consistent conditional fields (e.g. correct proportions of male/female smokers)."
    ])

    # ─────────────────────────────────────────────────────────────────────────
    # SLIDE 5: Pillar 2 - Clinical & Statistical Validation
    # ─────────────────────────────────────────────────────────────────────────
    add_bullet_slide(prs, "Pillar 2: Multi-Layer Validation Gates", [
        "1. Medical Rule Engine (Physiological Validation)",
        "  - Enforces clinical constraints on synthetic records.",
        "  - Examples: If Smoking Status = 'Never', Years Smoked must be 0; BMI is clipped between 15 and 45; Blood oxygen must remain within 80-100%.",
        "2. Statistical Distribution Matching",
        "  - Kolmogorov-Smirnov (KS) Test: Validates continuous features match target distributions.",
        "  - Chi-Square Test: Ensures categorical frequencies match baseline distributions.",
        "3. Data Drift Monitoring",
        "  - Population Stability Index (PSI): Ensures generated distributions do not shift from the baseline over time (Target PSI < 0.1)."
    ])

    # ─────────────────────────────────────────────────────────────────────────
    # SLIDE 6: Pillar 3 - Stacking Ensemble Classifier
    # ─────────────────────────────────────────────────────────────────────────
    add_bullet_slide(prs, "Pillar 3: Calibrated Stacking Ensemble", [
        "Base Learners (Level 0 Models):",
        "  - We train multiple diverse algorithms: XGBoost, CatBoost, Random Forest, and Gradient Boosting.",
        "  - Each model makes its own individual out-of-fold prediction probabilities.",
        "Meta-Learner (Level 1 Model):",
        "  - A Logistic Regression meta-classifier takes the prediction probabilities from Level 0 models as its features.",
        "  - It learns how to optimally weight the strengths of each model to produce a final uncalibrated risk score.",
        "Why Stacking?",
        "  - It reduces model variance and prevents overfitting, creating a robust predictor that generalizes well to unseen patient cases."
    ])

    # ─────────────────────────────────────────────────────────────────────────
    # SLIDE 7: Calibrated Stacking Ensemble Flowchart (NATIVE SHAPES)
    # ─────────────────────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(slide_layout)
    apply_background(slide, RGBColor(245, 248, 250))
    add_title(slide, "Calibrated Stacking Ensemble Architecture", "Flowchart Diagram of the Complete Prediction Pipeline")
    
    # Node Drawing Coordinates helper
    def draw_node(slide, text, left, top, width, height, bg_hex, text_color=RGBColor(0,0,0), border_color=RGBColor(180,180,180), font_size=11, bold=True):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor.from_string(bg_hex)
        shape.line.color.rgb = border_color
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0.05)
        p = tf.paragraphs[0]
        p.text = text
        p.alignment = PP_ALIGN.CENTER
        p.font.name = 'Calibri'
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = text_color
        return shape

    def draw_arrow(slide, x, y, length, vertical=True):
        # We can draw simple vertical/horizontal line arrows
        shape = slide.shapes.add_shape(
            MSO_SHAPE.DOWN_ARROW if vertical else MSO_SHAPE.RIGHT_ARROW, 
            Inches(x), Inches(y), 
            Inches(0.15) if vertical else Inches(length), 
            Inches(length) if vertical else Inches(0.15)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(120, 120, 120)
        shape.line.fill.background()

    def draw_label(slide, text, left, top, width, height):
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.alignment = PP_ALIGN.CENTER
        p.font.name = 'Calibri'
        p.font.size = Pt(9)
        p.font.color.rgb = RGBColor(100, 100, 100)

    # 1. Input Data
    draw_node(slide, "Input Patient Data", 5.16, 1.3, 3.0, 0.5, "E2EFDA", border_color=RGBColor(76,175,80))
    draw_arrow(slide, 6.58, 1.85, 0.35)

    # Yellow Box Layer 1: Base Learners
    container1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.2), Inches(2.2), Inches(8.9), Inches(1.3))
    container1.fill.solid()
    container1.fill.fore_color.rgb = RGBColor(255, 253, 230)
    container1.line.color.rgb = RGBColor(240, 220, 100)
    draw_label(slide, "Level 0 - Base Learners", 5.16, 2.22, 3.0, 0.25)
    
    # Inner Level 0 Models
    draw_node(slide, "XGBoost Classifier", 2.5, 2.6, 2.2, 0.7, "D9E1F2")
    draw_node(slide, "CatBoost Classifier", 5.5, 2.6, 2.2, 0.7, "D9E1F2")
    draw_node(slide, "Random Forest Classifier", 8.5, 2.6, 2.2, 0.7, "D9E1F2")
    
    # Arrows leaving level 0
    draw_arrow(slide, 3.5, 3.5, 0.3)
    draw_arrow(slide, 6.6, 3.5, 0.3)
    draw_arrow(slide, 9.5, 3.5, 0.3)
    
    draw_label(slide, "Probabilities", 2.9, 3.8, 1.2, 0.25)
    draw_label(slide, "Probabilities", 6.0, 3.8, 1.2, 0.25)
    draw_label(slide, "Probabilities", 8.9, 3.8, 1.2, 0.25)
    
    # Yellow Box Layer 2: Meta Learner
    container2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.66), Inches(4.1), Inches(6.0), Inches(1.0))
    container2.fill.solid()
    container2.fill.fore_color.rgb = RGBColor(255, 253, 230)
    container2.line.color.rgb = RGBColor(240, 220, 100)
    draw_label(slide, "Level 1 - Meta Learner", 5.16, 4.12, 3.0, 0.25)
    
    draw_node(slide, "Logistic Regression", 5.16, 4.45, 3.0, 0.5, "E2EFDA")
    
    draw_arrow(slide, 6.58, 5.15, 0.3)
    draw_label(slide, "Uncalibrated Probabilities", 5.16, 5.45, 3.0, 0.25)
    
    # Calibrator & Safety Threshold
    draw_node(slide, "Isotonic Regression Calibrator\n(Probability Calibration)", 3.16, 5.75, 3.2, 0.65, "FFF2CC", border_color=RGBColor(255,193,7))
    draw_arrow(slide, 6.58, 6.05, 1.5, vertical=False) # Horizontal arrow
    
    draw_node(slide, "Threshold Safety Optimizer\n(Decision Cutoff for 95% Recall)", 8.26, 5.75, 3.2, 0.65, "FFF2CC", border_color=RGBColor(255,193,7))
    
    draw_arrow(slide, 9.86, 6.45, 0.4)
    draw_node(slide, "Final Patient Risk Output", 8.26, 6.9, 3.2, 0.5, "E2EFDA", border_color=RGBColor(76,175,80))

    # ─────────────────────────────────────────────────────────────────────────
    # SLIDE 8: Probability Calibration & Threshold Optimization
    # ─────────────────────────────────────────────────────────────────────────
    add_bullet_slide(prs, "Probability Calibration & Safety Thresholds", [
        "Why Calibrate Prediction Scores?",
        "  - Raw machine learning probabilities are frequently skewed and uncalibrated.",
        "  - Isotonic Regression maps the raw stack outputs to true empirical percentages.",
        "  - If the model predicts an 85% risk, exactly 85% of such patients have lung cancer.",
        "Optimizing the Decision Cutoff:",
        "  - The standard classification threshold is 0.50.",
        "  - In medicine, false negatives are fatal. We must prioritize finding all actual cancer cases.",
        "  - We optimize the decision boundary (typically shifting to 0.30 - 0.35) to guarantee at least **95% Recall (Sensitivity)**.",
        "  - This ensures only a minimal percentage of positive cases are missed by the platform."
    ])

    # ─────────────────────────────────────────────────────────────────────────
    # SLIDE 9: Explainable AI (XAI) via SHAP
    # ─────────────────────────────────────────────────────────────────────────
    add_bullet_slide(prs, "Clinical Explainability via SHAP (XAI)", [
        "Demystifying the Prediction:",
        "  - Instead of giving a single number, the system utilizes **SHAP (SHapley Additive exPlanations)**.",
        "  - SHAP values are derived from game theory, computing the exact marginal contribution of each patient feature.",
        "Local Explanations (Patient-Level):",
        "  - Doctors can see exactly which factors pushed the patient's risk score higher or lower.",
        "  - Explanations are translated to simple clinical categories: 'Significantly Elevated Risk', 'Moderately Elevated Risk', or 'Protective Factor'.",
        "Example Output:",
        "  - *'Patient risk is 89.5%. Heaviest driver: Smoking history (+25% contribution), followed by high PM2.5 air exposure (+12% contribution).'*",
        "  - This builds clinical trust and makes the model easily auditable by doctors."
    ])

    # ─────────────────────────────────────────────────────────────────────────
    # SLIDE 10: Conclusion & Technical Highlights
    # ─────────────────────────────────────────────────────────────────────────
    add_bullet_slide(prs, "Conclusion & Project Highlights", [
        "✅ Solves Real-World Medical Constraints",
        "  - CTGAN generates privacy-compliant synthetic records, while statistical tests ensure distribution match.",
        "🩺 Clinical Validation & Safety Gates",
        "  - Medical rules prevent impossible records, while probability calibration ensures risk scores are reliable.",
        "🔍 Full Transparency (Explainable AI)",
        "  - Local explanations (SHAP) bridge the gap between machine learning and medical trust.",
        "📈 State-of-the-Art Architecture",
        "  - Combining modern generative models, statistical QA pipelines, stacking ensembles, and game-theoretic explainers.",
        "🎓 Perfect for Demo and Clinical Presentation!"
    ])
    
    prs.save("Lung_Cancer_Prediction_Presentation.pptx")
    print("Presentation created successfully as Lung_Cancer_Prediction_Presentation.pptx")

if __name__ == "__main__":
    create_presentation()
